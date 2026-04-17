#!/usr/bin/env python3
"""Generate slides for Multi-Chunk Posting design presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
import json

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT_BLUE = RGBColor(0x58, 0x9C, 0xF5)
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT_RED = RGBColor(0xF4, 0x71, 0x74)
ACCENT_YELLOW = RGBColor(0xE5, 0xC0, 0x7B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, bullets, left=1.0, top=2.0, width=11.0, font_size=20, color=WHITE, spacing=0.15):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, level, bcolor) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size - level * 2)
        p.font.color.rgb = bcolor or color
        p.font.name = 'Calibri'
        p.level = level
        p.space_after = Pt(spacing * 72)
        if level > 0:
            p.font.size = Pt(font_size - 2)
    return txBox

def add_title(slide, title, subtitle=None):
    add_text_box(slide, 0.8, 0.4, 11.5, 0.8, title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, 0.8, 1.1, 11.5, 0.5, subtitle, font_size=18, color=LIGHT_GRAY)

def add_code_box(slide, left, top, width, height, code, font_size=14):
    from pptx.util import Pt, Inches
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    # Add background
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x38)
    for i, line in enumerate(code.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN
        p.font.name = 'Consolas'
        p.space_after = Pt(2)
    return txBox

def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows)+1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    
    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    
    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x35, 0x35, 0x55)
    
    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r+1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = WHITE
                p.font.name = 'Calibri'
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            if r % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3A)
            else:
                cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x32)
    
    return table_shape

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)
add_text_box(slide, 1.0, 2.0, 11.0, 1.5, 
    "Multi-Chunk Posting for TiKV-backed SPTAG",
    font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.0, 3.5, 11.0, 0.8,
    "Optimizing SPFresh Insert Throughput with Write-Only Append",
    font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.0, 5.0, 11.0, 0.5,
    "SPTAG / SPFresh  •  April 2026",
    font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Problem Statement
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Problem: Read-Modify-Write Bottleneck")

add_bullet_slide(slide, [
    ("Traditional posting append in TiKV requires:", 0, WHITE),
    ("1. Get(key) → read full posting from TiKV", 1, ACCENT_RED),
    ("2. posting += new_vectors   (in-memory concatenation)", 1, ACCENT_RED),
    ("3. Put(key, posting) → write full posting back", 1, ACCENT_RED),
    ("", 0, WHITE),
    ("Problems at scale:", 0, WHITE),
    ("Posting sizes grow → Get latency grows linearly", 1, ACCENT_YELLOW),
    ("Write amplification: rewriting entire posting for small appends", 1, ACCENT_YELLOW),
    ("TiKV transaction overhead on large values", 1, ACCENT_YELLOW),
    ("Reassign queue explosion: ~20K jobs queued, only 16 threads", 1, ACCENT_YELLOW),
], top=1.8, font_size=22)

# ============================================================
# SLIDE 3: Solution Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Solution: Two Key Innovations")

add_bullet_slide(slide, [
    ("Innovation #1: Multi-Chunk Posting", 0, ACCENT_BLUE),
    ("Split each posting into multiple chunks with timestamp-suffixed keys", 1, WHITE),
    ("Append → Write-Only: just PutChunk with a new timestamp key (no Get needed)", 1, WHITE),
    ("Read → ScanPosting: prefix scan to collect all chunks for a head", 1, WHITE),
    ("Split → DeleteRange + write 2 new base chunks", 1, WHITE),
    ("", 0, WHITE),
    ("Innovation #2: Batch Reassign", 0, ACCENT_BLUE),
    ("Inline RNGSelection + batch Append by target head in CollectReAssign", 1, WHITE),
    ("Eliminates millions of queued async jobs → zero queue pressure", 1, WHITE),
    ("", 0, WHITE),
    ("Supporting Optimizations:", 0, LIGHT_GRAY),
    ("Eliminated posting cache (not needed with write-only append)", 1, LIGHT_GRAY),
    ("Increased thread pool (48) and TiKV gRPC stub pool (48)", 1, LIGHT_GRAY),
], top=1.8, font_size=22)

# ============================================================
# SLIDE 4: Key Format
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Key Format Design", "Leveraging TiKV's sorted key-value store")

add_code_box(slide, 1.0, 2.0, 11.0, 1.5,
    "Key Format:  [prefix]_[headID 4B]\\x00[timestamp 8B]\n"
    "\n"
    "Base chunk:   spfresh_sift1b_\\x00\\x00\\x03\\xE8\\x00    (headID=1000, no timestamp)\n"
    "Append chunk: spfresh_sift1b_\\x00\\x00\\x03\\xE8\\x00\\x17...  (headID=1000, ns timestamp)",
    font_size=16)

add_bullet_slide(slide, [
    ("\\x00 separator ensures base chunk sorts first (lowest key)", 0, WHITE),
    ("Nanosecond timestamp guarantees unique, ordered chunks", 0, WHITE),
    ("ScanPosting: prefix scan [headID\\x00, headID\\x01) collects all chunks", 0, WHITE),
    ("DeletePosting: single DeleteRange RPC removes all chunks for a head", 0, WHITE),
    ("PutBaseChunk: overwrites the base key (no timestamp suffix)", 0, WHITE),
], top=4.0, font_size=20)

# ============================================================
# SLIDE 5: Operation Mapping
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Operation Mapping")

headers = ["Operation", "Before (Single-Key)", "After (Multi-Chunk)", "TiKV RPCs"]
rows = [
    ["Append", "Get + Put (read-modify-write)", "PutChunk (write-only)", "2 → 1"],
    ["Read", "Get(key)", "ScanPosting (prefix scan)", "1 → 1"],
    ["Split", "Get + Delete + 2×Put", "ScanPosting + DeleteRange + 2×PutBase", "4 → 4"],
    ["Delete", "Delete(key)", "DeleteRange(prefix)", "1 → 1"],
    ["MultiGet", "MultiGet(keys)", "MultiScanPostings (parallel scans)", "1 → N (async)"],
]
add_table(slide, 0.5, 2.0, 12.3, 3.5, headers, rows, col_widths=[1.5, 3.5, 4.0, 1.5])

add_text_box(slide, 0.8, 5.8, 11.0, 0.5,
    "Key win: Append goes from 2 RPCs (Get+Put) to 1 RPC (PutChunk). This is the hot path.",
    font_size=18, color=ACCENT_GREEN)

# ============================================================
# SLIDE 6: Reassign Problem Deep-Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Innovation #2: The Reassign Bottleneck", "Discovered via custom statistics counters")

add_bullet_slide(slide, [
    ("What is Reassign?", 0, ACCENT_BLUE),
    ("After a posting is split into two new heads, vectors in nearby postings may now", 1, WHITE),
    ("be closer to the new heads. Reassign moves them to improve search quality.", 1, WHITE),
    ("", 0, WHITE),
    ("The hidden bottleneck (observed from runtime stats):", 0, ACCENT_RED),
    ("pending_queue: 2,241,000    split_inflight: 123,726    running: 16", 1, ACCENT_YELLOW),
    ("→ 2.1M items in queue were Reassign jobs, NOT Splits!", 1, ACCENT_RED),
    ("", 0, WHITE),
    ("Why so many?", 0, WHITE),
    ("Each Split triggers CollectReAssign with ReassignK=64", 1, WHITE),
    ("→ MultiGet 64 nearby heads, scan all their vectors", 1, WHITE),
    ("→ Hundreds of ReassignAsync jobs per Split, each requiring:", 1, WHITE),
    ("    1× RNGSelection (head index search) + 1× Append (TiKV RPC)", 2, ACCENT_YELLOW),
    ("16 threads processing millions of jobs = hours of queue drain time", 1, ACCENT_RED),
], top=1.8, font_size=20)

# ============================================================
# SLIDE 7: Batch Reassign - Before vs After
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Batch Reassign: Design Comparison")

# Before
add_text_box(slide, 0.5, 1.8, 5.8, 0.5, "Before: Per-Vector Async Jobs", font_size=22, color=ACCENT_RED, bold=True)
add_code_box(slide, 0.5, 2.4, 5.8, 3.2,
    "Split() completes → CollectReAssign()\n"
    "  for each vector needing reassign:\n"
    "    ReassignAsync(vector, headPrev)\n"
    "      // Creates ReassignAsyncJob\n"
    "      // Enqueues to splitThreadPool\n"
    "\n"
    "Each ReassignAsyncJob::exec():  // in pool\n"
    "  RNGSelection(vector) → targetHead\n"
    "  IncVersion(VID)\n"
    "  Append(targetHead, 1 vector) → PutChunk",
    font_size=13)

add_bullet_slide(slide, [
    ("N vectors → N jobs → N×RNGSelection + N×PutChunk RPCs", 0, ACCENT_RED),
    ("Jobs compete with Split/Merge for thread pool slots", 0, ACCENT_RED),
    ("Queue drains slowly: each job does network I/O", 0, ACCENT_RED),
], left=0.5, top=5.8, width=5.8, font_size=16)

# After
add_text_box(slide, 6.8, 1.8, 5.8, 0.5, "After: Inline Batch Processing", font_size=22, color=ACCENT_GREEN, bold=True)
add_code_box(slide, 6.8, 2.4, 5.8, 3.2,
    "Split() completes → CollectReAssign()\n"
    "  map<targetHead, string> batch;\n"
    "  for each vector needing reassign:\n"
    "    RNGSelection(vector) → targetHead\n"
    "    IncVersion(VID)\n"
    "    batch[targetHead] += vectorInfo\n"
    "\n"
    "  // Batch write phase:\n"
    "  for (targetHead, mergedData) in batch:\n"
    "    Append(targetHead, count, merged)\n"
    "    // Single PutChunk per target head",
    font_size=13)

add_bullet_slide(slide, [
    ("N vectors → ~K unique heads → K×PutChunk RPCs", 0, ACCENT_GREEN),
    ("No thread pool contention: runs inline in Split thread", 0, ACCENT_GREEN),
    ("Zero queue growth: Split returns fully completed", 0, ACCENT_GREEN),
], left=6.8, top=5.8, width=5.8, font_size=16)

# ============================================================
# SLIDE 8: Batch Reassign - Quantitative Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Batch Reassign: Quantitative Impact")

headers = ["Metric", "Before (Per-Vector Async)", "After (Inline Batch)"]
rows = [
    ["Jobs per Split", "~100-500 ReassignAsync jobs", "0 (inline)"],
    ["TiKV RPCs per Split", "N×PutChunk (one per vector)", "K×PutChunk (K = unique targets, K << N)"],
    ["Thread pool pressure", "Reassign jobs flood queue", "No queue impact"],
    ["Queue at Batch 2 (observed)", "2,241,000 pending items", "~0 Reassign items"],
    ["Effective parallelism", "16 threads (pool limit)", "48 threads (no queue bottleneck)"],
    ["RNGSelection", "Async, waits in queue", "Synchronous, immediate"],
    ["Batch 2 throughput", "~120 vec/s (queue-limited)", "~355 vec/s"],
]
add_table(slide, 0.3, 1.8, 12.7, 4.5, headers, rows, col_widths=[2.8, 4.5, 4.5])

add_bullet_slide(slide, [
    ("Key insight: RNGSelection is pure in-memory head index search (~microseconds).", 0, WHITE),
    ("The bottleneck was never the computation — it was the queuing and per-vector I/O.", 0, ACCENT_GREEN),
], top=6.0, font_size=18)

# ============================================================
# SLIDE 9: Deadlock Prevention
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Engineering Detail: Deadlock Prevention", "Discovered and fixed during batch reassign implementation")

add_code_box(slide, 0.8, 1.8, 11.0, 2.0,
    "// Potential recursive deadlock chain:\n"
    "Split(headA) → CollectReAssign() → batch Append(headB, threshold=3)\n"
    "  → Append detects overflow → synchronous Split(headB)\n"
    "    → CollectReAssign() → batch Append(headC) → ...\n"
    "// Same thread tries to re-acquire m_rwLocks → DEADLOCK",
    font_size=15)

add_bullet_slide(slide, [
    ("Problem: Batch Append with reassignThreshold > 0 triggers synchronous Split on overflow", 0, ACCENT_RED),
    ("→ Recursive call chain: Split → Reassign → Append → Split → ...", 1, ACCENT_RED),
    ("→ Same thread re-enters m_rwLocks mutex → 'Resource deadlock avoided' crash", 1, ACCENT_RED),
    ("", 0, WHITE),
    ("Solution: Batch Append uses reassignThreshold=0", 0, ACCENT_GREEN),
    ("→ On overflow, triggers SplitAsync (enqueues to thread pool) instead of sync Split", 1, ACCENT_GREEN),
    ("→ Breaks the recursive chain while maintaining correctness", 1, ACCENT_GREEN),
    ("", 0, WHITE),
    ("Original ReassignAsync design avoided this because each Reassign ran in a", 0, LIGHT_GRAY),
    ("separate thread pool job — the recursion was broken by the async boundary.", 1, LIGHT_GRAY),
], top=4.0, font_size=18)

# ============================================================
# SLIDE 10: Performance Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Performance Results: Insert Throughput", "SIFT1B, UInt8, dim=128, 1M vectors per batch, TiKV v9.0.0-beta.2")

headers = ["Batch", "Baseline (s)", "MultiChunk (s)", "Baseline (vec/s)", "MultiChunk (vec/s)", "Speedup"]
rows = [
    ["1", "3,253", "2,702", "307", "370", "1.2x"],
    ["2", "6,952", "2,821", "144", "355", "2.5x"],
    ["3", "8,143", "3,050", "123", "328", "2.7x"],
    ["4", "8,769", "3,204", "114", "312", "2.7x"],
    ["5", "9,146", "3,335", "109", "300", "2.7x"],
    ["6", "9,432", "3,415", "106", "293", "2.8x"],
    ["7", "9,577", "3,518", "104", "284", "2.7x"],
    ["8", "9,756", "3,598", "103", "278", "2.7x"],
    ["9", "9,761", "3,651", "102", "274", "2.7x"],
    ["10", "9,925", "3,671", "101", "272", "2.7x"],
]
add_table(slide, 0.3, 1.8, 12.7, 5.0, headers, rows, col_widths=[1.0, 2.0, 2.0, 2.2, 2.5, 1.2])

# ============================================================
# SLIDE 8: Search Quality
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Search Quality: Maintained", "Search latency and QPS remain comparable")

headers = ["Batch", "Baseline Latency (ms)", "MultiChunk Latency (ms)", "Baseline QPS", "MultiChunk QPS"]
rows = [
    ["1", "11.58", "10.38", "342.9", "382.9"],
    ["2", "12.22", "11.34", "324.9", "350.4"],
    ["3", "12.40", "11.25", "320.1", "353.0"],
    ["4", "12.48", "12.26", "318.9", "324.4"],
    ["5", "13.02", "12.33", "305.6", "322.1"],
    ["6", "13.11", "13.33", "303.5", "298.8"],
    ["7", "13.40", "13.33", "296.7", "298.2"],
    ["8", "13.83", "13.75", "287.6", "289.1"],
    ["9", "13.87", "14.20", "287.2", "280.1"],
    ["10", "14.44", "14.49", "275.8", "274.8"],
]
add_table(slide, 0.8, 1.8, 11.5, 5.0, headers, rows, col_widths=[1.2, 2.5, 2.8, 2.0, 2.2])

add_text_box(slide, 0.8, 6.5, 11.0, 0.5,
    "Search latency is comparable or slightly better in early batches. No quality regression.",
    font_size=16, color=ACCENT_GREEN)

# ============================================================
# SLIDE 9: Key Takeaways
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Key Takeaways")

add_bullet_slide(slide, [
    ("Two complementary innovations deliver 2.5-2.8x throughput improvement", 0, ACCENT_GREEN),
    ("", 0, WHITE),
    ("#1 Multi-Chunk Posting: write-only append eliminates read-modify-write", 0, ACCENT_BLUE),
    ("Reduces Append from 2 TiKV RPCs to 1 (the hot path)", 1, WHITE),
    ("", 0, WHITE),
    ("#2 Batch Reassign: inline processing eliminates queue explosion", 0, ACCENT_BLUE),
    ("Reduces N async jobs to K batch writes (K << N)", 1, WHITE),
    ("Prevents baseline degradation from 307 → 100 vec/s", 1, ACCENT_RED),
    ("", 0, WHITE),
    ("Throughput stability: 270-370 vec/s maintained across 10 batches", 0, WHITE),
    ("Search quality preserved: latency and QPS comparable to baseline", 0, WHITE),
    ("Design is generic: applicable to any KV store with prefix scan support", 0, LIGHT_GRAY),
], top=1.8, font_size=24)

# ============================================================
# SLIDE 10: Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Configuration & Setup")

add_bullet_slide(slide, [
    ("Hardware: 48-core CPU, 378GB RAM, NVMe striped storage", 0, WHITE),
    ("TiKV: v9.0.0-beta.2 nightly, 3 TiKV + 3 PD nodes (Docker)", 0, WHITE),
    ("    grpc-concurrency=16, block-cache=60GB, region-max-size=512MB", 1, LIGHT_GRAY),
    ("", 0, WHITE),
    ("SPTAG Config:", 0, ACCENT_BLUE),
    ("NumInsertThreads=16, AppendThreadNum=48", 1, WHITE),
    ("kStubPoolSize=48 (TiKV gRPC connection pool)", 1, WHITE),
    ("UseMultiChunkPosting=true", 1, WHITE),
    ("ReassignK=64 (from indexloader.ini defaults)", 1, WHITE),
    ("", 0, WHITE),
    ("Dataset: SIFT1B (UInt8, dim=128, 1B vectors), L2 distance", 0, WHITE),
    ("Benchmark: 999 batches × 1M vectors, 200 queries per batch", 0, WHITE),
], top=1.8, font_size=20)

# Save
output_path = '/home/azureuser/qiazh/SPTAG/evaluation/MultiChunk_Posting_Design.pptx'
prs.save(output_path)
print(f"Saved to {output_path}")
